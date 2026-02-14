"""PowerPoint adapter for PresentationTool."""

from __future__ import annotations

import asyncio
import io
from collections.abc import Sequence
from typing import Any

from fireflyframework_genai.tools.base import GuardProtocol

from firefly_dworkers.tools.presentation.base import PresentationTool
from firefly_dworkers.tools.presentation.models import (
    PlaceholderInfo,
    PresentationData,
    SlideData,
    SlideOperation,
    SlideSpec,
)
from firefly_dworkers.tools.registry import tool_registry

try:
    import pptx
    from lxml import etree

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ── Table styling helpers (module-level, private) ──────────────────────────

_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _hex_to_srgb(hex_color: str) -> str:
    """Strip '#' and return uppercase 6-char hex for XML ``srgbClr``."""
    return hex_color.lstrip("#").upper()


def _set_cell_fill(cell: Any, hex_color: str) -> None:
    """Set solid fill on a table cell via XML manipulation."""
    tc_pr = cell._tc.get_or_add_tcPr()
    # Remove any existing fill
    for old_fill in tc_pr.findall(f"{{{_DRAWINGML_NS}}}solidFill"):
        tc_pr.remove(old_fill)
    solid_fill = etree.SubElement(tc_pr, f"{{{_DRAWINGML_NS}}}solidFill")
    etree.SubElement(solid_fill, f"{{{_DRAWINGML_NS}}}srgbClr", val=_hex_to_srgb(hex_color))


def _set_cell_borders(cell: Any, hex_color: str, width_pt: float = 0.5) -> None:
    """Set thin borders on all four sides of a table cell via XML."""
    width_emu = int(width_pt * 12700)  # 1 pt = 12700 EMU
    tc_pr = cell._tc.get_or_add_tcPr()
    for side in ("lnL", "lnR", "lnT", "lnB"):
        # Remove existing
        for old in tc_pr.findall(f"{{{_DRAWINGML_NS}}}{side}"):
            tc_pr.remove(old)
        ln = etree.SubElement(tc_pr, f"{{{_DRAWINGML_NS}}}{side}", w=str(width_emu))
        solid_fill = etree.SubElement(ln, f"{{{_DRAWINGML_NS}}}solidFill")
        etree.SubElement(solid_fill, f"{{{_DRAWINGML_NS}}}srgbClr", val=_hex_to_srgb(hex_color))


def _set_cell_margins(
    cell: Any,
    left: int = 91440,
    top: int = 45720,
    right: int = 91440,
    bottom: int = 45720,
) -> None:
    """Set cell padding via tcPr margin attributes (values in EMU)."""
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_pr.set("marL", str(left))
    tc_pr.set("marT", str(top))
    tc_pr.set("marR", str(right))
    tc_pr.set("marB", str(bottom))


def _set_run_font_color(run: Any, hex_color: str) -> None:
    """Set the font color on a python-pptx run from a hex string."""
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    run.font.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _require_pptx() -> None:
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx required: pip install firefly-dworkers[presentation]")


@tool_registry.register("powerpoint", category="presentation")
class PowerPointTool(PresentationTool):
    """Read, create, and modify PowerPoint presentations using python-pptx."""

    def __init__(
        self,
        *,
        timeout: float = 60.0,
        guards: Sequence[GuardProtocol] = (),
    ) -> None:
        super().__init__(
            "powerpoint",
            description="Read, create, and modify PowerPoint (.pptx) presentations.",
            timeout=timeout,
            guards=guards,
        )

    async def _read_presentation(self, source: str) -> PresentationData:
        _require_pptx()
        return await asyncio.to_thread(self._read_sync, source)

    async def _create_presentation(self, template: str, slides: list[SlideSpec]) -> bytes:
        _require_pptx()
        return await asyncio.to_thread(self._create_sync, template, slides)

    async def _modify_presentation(self, source: str, operations: list[SlideOperation]) -> bytes:
        _require_pptx()
        return await asyncio.to_thread(self._modify_sync, source, operations)

    # -- Sync implementations --

    def _read_sync(self, source: str) -> PresentationData:
        prs = pptx.Presentation(source)
        slides = []
        for i, slide in enumerate(prs.slides):
            title = ""
            content_parts: list[str] = []
            placeholders: list[PlaceholderInfo] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.placeholder_format is not None and shape.placeholder_format.idx == 0:
                        title = shape.text_frame.text
                    else:
                        content_parts.append(shape.text_frame.text)

                if shape.placeholder_format is not None:
                    placeholders.append(
                        PlaceholderInfo(
                            idx=shape.placeholder_format.idx,
                            name=shape.name,
                        )
                    )

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text

            slides.append(
                SlideData(
                    index=i,
                    layout=slide.slide_layout.name,
                    title=title,
                    content="\n".join(content_parts),
                    placeholders=placeholders,
                    notes=notes,
                )
            )

        layouts = [layout.name for layout in prs.slide_layouts]
        return PresentationData(
            slides=slides,
            master_layouts=layouts,
            slide_width=prs.slide_width or 0,
            slide_height=prs.slide_height or 0,
        )

    def _create_sync(self, template: str, slides: list[SlideSpec]) -> bytes:
        prs = pptx.Presentation(template) if template else pptx.Presentation()

        # Remove all existing slides from the template so we start fresh
        # while preserving masters, layouts, and theme assets.
        if template:
            self._remove_existing_slides(prs)

        for spec in slides:
            layout = self._find_layout(prs, spec.layout)
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title and spec.title:
                slide.shapes.title.text = spec.title
                if spec.title_style:
                    self._apply_text_style(slide.shapes.title.text_frame, spec.title_style)

            body_ph = self._find_body_placeholder(slide)
            if body_ph:
                if spec.bullet_points:
                    from pptx.util import Pt as _Pt

                    tf = body_ph.text_frame
                    tf.clear()
                    for j, point in enumerate(spec.bullet_points):
                        if j == 0:
                            tf.text = point
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                            p.text = point
                        p.level = 0
                        p.space_before = _Pt(4)
                        p.space_after = _Pt(2)
                elif spec.content:
                    body_ph.text_frame.text = spec.content

                if spec.body_style:
                    self._apply_text_style(body_ph.text_frame, spec.body_style)

            if spec.table:
                self._add_table(slide, spec.table)

            if spec.chart:
                self._add_chart(slide, spec.chart)

            if spec.image_path:
                self._add_image(slide, spec.image_path)

            for img in spec.images:
                if img.file_path:
                    self._add_image_placement(slide, img)

            if spec.speaker_notes:
                slide.notes_slide.notes_text_frame.text = spec.speaker_notes

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _modify_sync(self, source: str, operations: list[SlideOperation]) -> bytes:
        prs = pptx.Presentation(source)
        for op in operations:
            if op.operation == "update_content":
                slide = prs.slides[op.slide_index]
                if "title" in op.data and slide.shapes.title:
                    slide.shapes.title.text = op.data["title"]
            elif op.operation == "add_slide":
                layout = self._find_layout(prs, op.data.get("layout", "Blank"))
                prs.slides.add_slide(layout)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # -- Helpers --

    @staticmethod
    def _remove_existing_slides(prs: Any) -> None:
        """Remove all existing slides from a presentation, keeping masters/layouts."""
        sldIdLst = prs._element.sldIdLst
        if sldIdLst is None:
            return
        slide_ids = list(sldIdLst)
        for sldId in slide_ids:
            rId = sldId.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            sldIdLst.remove(sldId)
            if rId:
                try:
                    prs.part.rels.pop(rId)
                except KeyError:
                    pass

    @staticmethod
    def _all_layouts(prs: Any) -> list[Any]:
        """Collect layouts from ALL slide masters (not just the first)."""
        layouts: list[Any] = []
        seen: set[str] = set()
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name not in seen:
                    seen.add(layout.name)
                    layouts.append(layout)
        return layouts

    @classmethod
    def _find_layout(cls, prs: Any, name: str) -> Any:
        """Find a layout by name with fuzzy/semantic matching.

        Searches ALL slide masters. Tries exact match first, then
        case-insensitive, then keyword-based mapping for common
        English / Spanish layout names.
        """
        all_layouts = cls._all_layouts(prs)

        # Exact match
        for layout in all_layouts:
            if layout.name == name:
                return layout

        # Case-insensitive match
        name_lower = name.lower()
        for layout in all_layouts:
            if layout.name.lower() == name_lower:
                return layout

        # Semantic keyword mapping: English concept → keywords to match
        _CONCEPT_KEYWORDS: dict[str, list[str]] = {
            "title slide": ["portada", "cover"],
            "title and content": ["título de 2", "título de 1"],
            "two content": ["dos", "two content", "2 columnas"],
            "section header": ["fondo", "sección", "section", "divider"],
            "blank": ["blanco", "blank"],
            "back cover": ["contraportada", "back"],
        }

        # Find the concept for the requested name
        for concept, keywords in _CONCEPT_KEYWORDS.items():
            if concept in name_lower or name_lower in concept:
                for layout in all_layouts:
                    layout_lower = layout.name.lower()
                    for kw in keywords:
                        if kw in layout_lower:
                            return layout
                break

        # Reverse: if the requested name itself contains keywords
        for layout in all_layouts:
            layout_lower = layout.name.lower()
            for word in name_lower.split():
                if len(word) > 3 and word in layout_lower:
                    return layout

        # Best-effort: pick a content layout (with ph[0] title + ph[13] body)
        for layout in all_layouts:
            ph_indices = {ph.placeholder_format.idx for ph in layout.placeholders}
            if 0 in ph_indices and 13 in ph_indices:
                return layout

        # Fallback: any layout with a title placeholder
        for layout in all_layouts:
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:
                    return layout

        return prs.slide_layouts[0]

    @staticmethod
    def _find_body_placeholder(slide: Any) -> Any:
        """Find the body/content placeholder in a slide.

        Tries standard idx 1 first, then idx 13 (common in custom templates),
        then any non-title text placeholder.
        """
        for target_idx in (1, 13, 14):
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == target_idx:
                    return shape
        # Fallback: any non-title placeholder with a text frame
        for shape in slide.placeholders:
            if shape.placeholder_format.idx != 0 and shape.has_text_frame:
                return shape
        return None

    @staticmethod
    def _add_table(slide: Any, table_spec: Any) -> None:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        if hasattr(table_spec, "headers"):
            headers = table_spec.headers
            rows_data = table_spec.rows
        else:
            headers = table_spec.get("headers", [])
            rows_data = table_spec.get("rows", [])

        n_rows = len(rows_data) + 1
        n_cols = len(headers)
        if n_cols == 0:
            return

        # Extract styling fields with fallback defaults
        _attr = lambda name, default: getattr(table_spec, name, None) or (table_spec.get(name) if isinstance(table_spec, dict) else None) or default  # noqa: E731
        header_bg = _attr("header_bg_color", "")
        header_text_color = _attr("header_text_color", "#FFFFFF")
        alternating = _attr("alternating_rows", True)
        alt_color = _attr("alt_row_color", "#F5F5F5")
        border_color = _attr("border_color", "#CCCCCC")
        font_name = _attr("font_name", "")
        header_font_size = float(_attr("header_font_size", 10.0))
        cell_font_size = float(_attr("cell_font_size", 9.0))

        row_height = Inches(0.35)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(1), Inches(2), Inches(8), row_height * n_rows
        )
        table = table_shape.table

        # Disable built-in banding so we control row colors manually
        tbl = table._tbl
        tbl_pr = tbl.tblPr
        if tbl_pr is not None:
            tbl_pr.set("bandRow", "0")
            tbl_pr.set("bandCol", "0")
            tbl_pr.set("firstRow", "0")
            tbl_pr.set("lastRow", "0")

        # ── Header row ──
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            for run in tf.paragraphs[0].runs:
                run.font.bold = True
                run.font.size = Pt(header_font_size)
                if font_name:
                    run.font.name = font_name
                if header_bg:
                    _set_run_font_color(run, header_text_color)

            if header_bg:
                _set_cell_fill(cell, header_bg)
            _set_cell_borders(cell, border_color)
            _set_cell_margins(cell)

        # ── Data rows ──
        for r, row in enumerate(rows_data):
            is_even = r % 2 == 0
            for c, val in enumerate(row):
                if c >= n_cols:
                    continue
                cell = table.cell(r + 1, c)
                cell.text = str(val)
                tf = cell.text_frame
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                for run in tf.paragraphs[0].runs:
                    run.font.size = Pt(cell_font_size)
                    if font_name:
                        run.font.name = font_name
                    _set_run_font_color(run, "#333333")

                # Alternating row fill
                if alternating and is_even:
                    _set_cell_fill(cell, alt_color)
                _set_cell_borders(cell, border_color)
                _set_cell_margins(cell)

    @staticmethod
    def _add_chart(slide: Any, chart_spec: Any) -> None:
        """Add a native chart to a slide using ChartRenderer."""
        from firefly_dworkers.design.charts import ChartRenderer
        from firefly_dworkers.design.models import DataSeries, ResolvedChart

        renderer = ChartRenderer()
        resolved = ResolvedChart(
            chart_type=chart_spec.chart_type,
            title=chart_spec.title,
            categories=chart_spec.categories,
            series=[DataSeries(name=s["name"], values=s["values"]) for s in chart_spec.series],
            colors=chart_spec.colors,
            show_legend=chart_spec.show_legend,
            show_data_labels=chart_spec.show_data_labels,
            stacked=chart_spec.stacked,
        )
        renderer.render_for_pptx(resolved, slide)

    @staticmethod
    def _add_image(slide: Any, image_path: str) -> None:
        """Add a single image from a file path to a slide."""
        from pptx.util import Inches

        slide.shapes.add_picture(image_path, Inches(1), Inches(2), Inches(4), Inches(3))

    @staticmethod
    def _add_image_placement(slide: Any, img: Any) -> None:
        """Add an image with explicit placement from an ImagePlacement spec."""
        from pptx.util import Emu, Inches

        slide.shapes.add_picture(
            img.file_path,
            Emu(int(img.left)) if img.left else Inches(1),
            Emu(int(img.top)) if img.top else Inches(2),
            Emu(int(img.width)) if img.width else None,
            Emu(int(img.height)) if img.height else None,
        )

    @staticmethod
    def _apply_text_style(text_frame: Any, style: Any) -> None:
        """Apply TextStyle to all runs in a text frame.

        Handles both paragraph-level properties (alignment) and
        run-level properties (font, size, bold, italic, color).
        """
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        if style is None:
            return

        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }

        for paragraph in text_frame.paragraphs:
            # Paragraph-level: alignment
            if hasattr(style, "alignment") and style.alignment in alignment_map:
                paragraph.alignment = alignment_map[style.alignment]

            # Run-level: font properties
            for run in paragraph.runs:
                if style.font_name:
                    run.font.name = style.font_name
                if style.font_size > 0:
                    run.font.size = Pt(style.font_size)
                if style.bold:
                    run.font.bold = True
                if style.italic:
                    run.font.italic = True
                if style.color:
                    hex_color = style.color.lstrip("#")
                    run.font.color.rgb = RGBColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16),
                        int(hex_color[4:6], 16),
                    )
