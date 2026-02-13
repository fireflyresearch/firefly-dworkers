"""ChartRenderer -- renders ResolvedChart to format-specific chart objects.

Produces native charts for PPTX (python-pptx) and XLSX (openpyxl),
and rasterised PNG images via matplotlib for DOCX/PDF embedding.
"""

from __future__ import annotations

import asyncio
import io
import logging
from typing import Any

from firefly_dworkers.design.models import ResolvedChart

# ── Lazy library imports ────────────────────────────────────────────────────

try:
    import matplotlib

    matplotlib.use("Agg")  # Non-interactive backend — must precede pyplot
    import matplotlib.pyplot as plt

    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False

try:
    import pptx
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    from openpyxl.chart import (
        AreaChart,
        BarChart,
        DoughnutChart,
        LineChart,
        PieChart,
        Reference,
        ScatterChart,
    )

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

logger = logging.getLogger(__name__)

# ── Chart-type mappings ─────────────────────────────────────────────────────

_PPTX_CHART_TYPES: dict[str, Any] = {}
if PPTX_AVAILABLE:
    _PPTX_CHART_TYPES = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "area": XL_CHART_TYPE.AREA,
        "area_stacked": XL_CHART_TYPE.AREA_STACKED,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }

_XLSX_CHART_CLASSES: dict[str, type] = {}
if OPENPYXL_AVAILABLE:
    _XLSX_CHART_CLASSES = {
        "bar": BarChart,
        "line": LineChart,
        "pie": PieChart,
        "scatter": ScatterChart,
        "area": AreaChart,
        "doughnut": DoughnutChart,
    }


class ChartRenderer:
    """Renders :class:`ResolvedChart` to format-specific chart objects."""

    # ── matplotlib (PNG bytes) ──────────────────────────────────────────────

    async def render_to_image(
        self,
        chart: ResolvedChart,
        *,
        width: int = 800,
        height: int = 600,
    ) -> bytes:
        """Render *chart* to PNG bytes using matplotlib.

        Used for DOCX / PDF embedding where native chart objects are not
        available.  The heavy matplotlib work runs in a thread via
        :func:`asyncio.to_thread` to avoid blocking the event loop.
        """
        if not MATPLOTLIB_AVAILABLE:
            raise ImportError("matplotlib required: pip install matplotlib")
        return await asyncio.to_thread(self.render_to_image_sync, chart, width, height)

    def render_to_image_sync(self, chart: ResolvedChart, width: int = 800, height: int = 600) -> bytes:
        """Synchronous variant of :meth:`render_to_image`.

        Useful when called from within a sync context that is already
        off the event loop (e.g. ``_create_sync`` in tool adapters).
        """
        fig, ax = plt.subplots(figsize=(width / 100, height / 100), dpi=100)
        try:
            self._plot_matplotlib(chart, ax)
            buf = io.BytesIO()
            fig.savefig(buf, format="png", bbox_inches="tight", dpi=100)
            return buf.getvalue()
        finally:
            plt.close(fig)

    def _plot_matplotlib(self, chart: ResolvedChart, ax: Any) -> None:
        """Dispatch to the correct matplotlib plot method."""
        categories = chart.categories or []
        colors = chart.colors or None
        chart_type = chart.chart_type.lower()

        if chart_type == "bar":
            self._plot_bar(chart, ax, categories, colors)
        elif chart_type == "line":
            self._plot_line(chart, ax, categories, colors)
        elif chart_type == "pie":
            self._plot_pie(chart, ax, categories, colors)
        elif chart_type == "scatter":
            self._plot_scatter(chart, ax, colors)
        elif chart_type == "area":
            self._plot_area(chart, ax, categories, colors)
        elif chart_type == "doughnut":
            self._plot_doughnut(chart, ax, categories, colors)
        else:
            # Best-effort fallback: treat unknown types as bar charts
            logger.warning("Unsupported chart type '%s', falling back to bar chart", chart_type)
            self._plot_bar(chart, ax, categories, colors)

        if chart.title:
            ax.set_title(chart.title)
        if chart.show_legend and chart_type not in ("pie", "doughnut"):
            ax.legend()

    # -- matplotlib helpers ---------------------------------------------------

    @staticmethod
    def _plot_bar(chart: ResolvedChart, ax: Any, categories: list[str], colors: list[str] | None) -> None:
        import numpy as np

        x = np.arange(len(categories))
        n_series = len(chart.series)
        bar_width = 0.8 / max(n_series, 1)

        bottom = np.zeros(len(categories)) if chart.stacked else None

        for idx, s in enumerate(chart.series):
            values = [float(v) for v in s.values]
            color = colors[idx % len(colors)] if colors else None
            if chart.stacked:
                ax.bar(x, values, bar_width * n_series, bottom=bottom, label=s.name, color=color)
                bottom = bottom + np.array(values)  # type: ignore[operator]
            else:
                ax.bar(x + idx * bar_width, values, bar_width, label=s.name, color=color)

        ax.set_xticks(x + bar_width * (n_series - 1) / 2 if not chart.stacked else x)
        ax.set_xticklabels(categories)

        if chart.show_data_labels:
            for container in ax.containers:
                ax.bar_label(container)

    @staticmethod
    def _plot_line(chart: ResolvedChart, ax: Any, categories: list[str], colors: list[str] | None) -> None:
        for idx, s in enumerate(chart.series):
            values = [float(v) for v in s.values]
            color = colors[idx % len(colors)] if colors else None
            ax.plot(categories, values, label=s.name, color=color, marker="o")

        if chart.show_data_labels:
            for s in chart.series:
                for i, v in enumerate(s.values):
                    ax.annotate(str(v), (categories[i], float(v)))

    @staticmethod
    def _plot_pie(chart: ResolvedChart, ax: Any, categories: list[str], colors: list[str] | None) -> None:
        # Pie chart uses only the first series
        if not chart.series:
            return
        values = [float(v) for v in chart.series[0].values]
        kwargs: dict[str, Any] = {"labels": categories, "autopct": "%1.1f%%" if chart.show_data_labels else None}
        if colors:
            kwargs["colors"] = colors[: len(values)]
        ax.pie(values, **kwargs)

    @staticmethod
    def _plot_scatter(chart: ResolvedChart, ax: Any, colors: list[str] | None) -> None:
        # For scatter, each series provides (x, y) pairs; categories as x if provided
        for idx, s in enumerate(chart.series):
            values = [float(v) for v in s.values]
            x = list(range(len(values)))
            color = colors[idx % len(colors)] if colors else None
            ax.scatter(x, values, label=s.name, color=color)

    @staticmethod
    def _plot_area(chart: ResolvedChart, ax: Any, categories: list[str], colors: list[str] | None) -> None:
        for idx, s in enumerate(chart.series):
            values = [float(v) for v in s.values]
            color = colors[idx % len(colors)] if colors else None
            ax.fill_between(range(len(values)), values, alpha=0.5, label=s.name, color=color)
        if categories:
            ax.set_xticks(range(len(categories)))
            ax.set_xticklabels(categories)

    @staticmethod
    def _plot_doughnut(chart: ResolvedChart, ax: Any, categories: list[str], colors: list[str] | None) -> None:
        if not chart.series:
            return
        values = [float(v) for v in chart.series[0].values]
        kwargs: dict[str, Any] = {
            "labels": categories,
            "wedgeprops": {"width": 0.4},
            "autopct": "%1.1f%%" if chart.show_data_labels else None,
        }
        if colors:
            kwargs["colors"] = colors[: len(values)]
        ax.pie(values, **kwargs)

    # ── python-pptx (native slide chart) ────────────────────────────────────

    def render_for_pptx(
        self,
        chart: ResolvedChart,
        slide: Any,
        *,
        left: float = 914400,
        top: float = 1524000,
        width: float = 7315200,
        height: float = 4572000,
    ) -> None:
        """Add a native chart to a python-pptx *slide*.

        Positional arguments are in EMU (English Metric Units).
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx required: pip install firefly-dworkers[presentation]")

        chart_type = chart.chart_type.lower()
        is_scatter = chart_type == "scatter"

        # Resolve XL_CHART_TYPE enum member
        if chart.stacked and chart_type in ("bar", "area"):
            xl_type = _PPTX_CHART_TYPES.get(f"{chart_type}_stacked")
        else:
            xl_type = _PPTX_CHART_TYPES.get(chart_type)

        if xl_type is None:
            logger.warning("Unsupported pptx chart type '%s', falling back to COLUMN_CLUSTERED", chart_type)
            xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        # Build chart data
        if is_scatter:
            chart_data = XyChartData()
            for s in chart.series:
                xy_series = chart_data.add_series(s.name)
                for i, v in enumerate(s.values):
                    xy_series.add_data_point(i, float(v))
        else:
            chart_data = CategoryChartData()
            chart_data.categories = chart.categories
            for s in chart.series:
                chart_data.add_series(s.name, [float(v) for v in s.values])

        chart_shape = slide.shapes.add_chart(
            xl_type,
            Emu(int(left)),
            Emu(int(top)),
            Emu(int(width)),
            Emu(int(height)),
            chart_data,
        )
        chart_obj = chart_shape.chart
        chart_obj.has_legend = chart.show_legend

        if chart.title:
            chart_obj.has_title = True
            chart_obj.chart_title.text_frame.text = chart.title

    # ── openpyxl (native worksheet chart) ────────────────────────────────────

    def render_for_xlsx(
        self,
        chart: ResolvedChart,
        ws: Any,
        *,
        cell: str = "E1",
    ) -> None:
        """Add a native chart to an openpyxl worksheet.

        Data is written to a hidden helper area starting at column
        ``_data_start_col`` so the chart's :class:`Reference` objects can
        point at real cells.
        """
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl required: pip install firefly-dworkers[data]")

        chart_type = chart.chart_type.lower()

        # Write data to the worksheet so chart Reference objects can use it.
        # Use a far-right column area to avoid collisions with user content.
        data_start_col = 50  # Column AX — far enough to not clash
        data_start_row = 1

        # Row 1: header (empty first cell + category labels)
        ws.cell(row=data_start_row, column=data_start_col, value="")
        for ci, cat in enumerate(chart.categories):
            ws.cell(row=data_start_row, column=data_start_col + 1 + ci, value=cat)

        # Row 2+: series name in first cell, then values
        for si, s in enumerate(chart.series):
            row = data_start_row + 1 + si
            ws.cell(row=row, column=data_start_col, value=s.name)
            for vi, v in enumerate(s.values):
                ws.cell(row=row, column=data_start_col + 1 + vi, value=float(v))

        # Determine the chart class
        chart_cls = _XLSX_CHART_CLASSES.get(chart_type)
        if chart_cls is None:
            logger.warning("Unsupported xlsx chart type '%s', falling back to BarChart", chart_type)
            chart_cls = BarChart

        chart_obj = chart_cls()

        if chart.title:
            chart_obj.title = chart.title

        # Category reference (first row, excluding the header cell)
        n_categories = len(chart.categories)
        n_series = len(chart.series)

        cats = Reference(
            ws,
            min_col=data_start_col + 1,
            max_col=data_start_col + n_categories,
            min_row=data_start_row,
        )

        for si in range(n_series):
            values_ref = Reference(
                ws,
                min_col=data_start_col + 1,
                max_col=data_start_col + n_categories,
                min_row=data_start_row + 1 + si,
            )
            chart_obj.add_data(values_ref, from_rows=True, titles_from_data=False)

        chart_obj.set_categories(cats)

        # Series naming
        for si, s in enumerate(chart.series):
            if si < len(chart_obj.series):
                chart_obj.series[si].title = openpyxl.chart.series.SeriesLabel(v=s.name)  # type: ignore[attr-defined]

        # Stacked grouping for bar/area
        if chart.stacked and chart_type in ("bar", "area"):
            chart_obj.grouping = "stacked"

        # Legend
        if not chart.show_legend:
            chart_obj.legend = None

        ws.add_chart(chart_obj, cell)
