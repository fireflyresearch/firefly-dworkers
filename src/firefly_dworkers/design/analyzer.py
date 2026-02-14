"""TemplateAnalyzer — extracts design DNA from existing documents.

Analyzes PPTX, DOCX, and XLSX files to produce a :class:`DesignProfile`
capturing colors, fonts, layouts, and other design attributes.
"""

from __future__ import annotations

import asyncio
import os
from typing import Any
from xml.etree.ElementTree import Element

from firefly_dworkers.design.models import DesignProfile, LayoutZone, PlaceholderZone

# ── Lazy library imports ────────────────────────────────────────────────────

try:
    import pptx

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False


# ── Constants ───────────────────────────────────────────────────────────────

_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_EMU_PER_INCH = 914400

_SUPPORTED_FORMATS: dict[str, str] = {
    ".pptx": "pptx",
    ".docx": "docx",
    ".xlsx": "xlsx",
}


class TemplateAnalyzer:
    """Analyzes existing documents to extract their design language (colors, fonts, layouts)."""

    def __init__(self, *, vlm_model: str = "") -> None:
        self._vlm_model = vlm_model

    # ── Public API ──────────────────────────────────────────────────────────

    def _detect_format(self, source: str) -> str:
        """Detect format from file extension.

        Returns ``'pptx'``, ``'docx'``, or ``'xlsx'``.

        Raises:
            ValueError: For unsupported or unrecognised extensions.
        """
        _, ext = os.path.splitext(source)
        ext = ext.lower()
        fmt = _SUPPORTED_FORMATS.get(ext)
        if fmt is None:
            raise ValueError(f"Unsupported format: '{ext}'. Expected one of {sorted(_SUPPORTED_FORMATS.keys())}")
        return fmt

    async def analyze(self, source: str) -> DesignProfile:
        """Auto-detect format and extract design DNA into a :class:`DesignProfile`."""
        fmt = self._detect_format(source)
        dispatch = {
            "pptx": self._analyze_pptx,
            "docx": self._analyze_docx,
            "xlsx": self._analyze_xlsx,
        }
        profile = await dispatch[fmt](source)

        if self._vlm_model and self._is_profile_empty(profile) and fmt == "pptx":
            try:
                profile = await self._vlm_analyze(source, profile)
            except Exception:
                pass  # XML result stands

        return profile

    # ── PPTX analysis ──────────────────────────────────────────────────────

    async def _analyze_pptx(self, path: str) -> DesignProfile:
        """Extract from PPTX: layout names, theme colors, font scheme, slide dimensions."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx required: pip install firefly-dworkers[presentation]")
        profile = await asyncio.to_thread(self._analyze_pptx_sync, path)

        # LLM-based classification for custom placeholders
        if self._vlm_model and profile.layout_zones:
            try:
                profile.layout_zones = await self._classify_zones_with_llm(
                    profile.layout_zones
                )
            except Exception:
                pass  # Spec-based classification stands

        return profile

    def _analyze_pptx_sync(self, path: str) -> DesignProfile:
        prs = pptx.Presentation(path)

        # Layouts — collect from ALL slide masters, not just the first
        seen: set[str] = set()
        available_layouts: list[str] = []
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name not in seen:
                    seen.add(layout.name)
                    available_layouts.append(layout.name)

        # Theme colors from XML
        color_palette: list[str] = []
        primary_color = ""
        secondary_color = ""
        accent_color = ""
        try:
            master_element = prs.slide_masters[0].element
            color_palette, primary_color, secondary_color, accent_color = self._extract_theme_colors(master_element)
        except Exception:
            pass

        # Font scheme from XML
        heading_font = ""
        body_font = ""
        try:
            master_element = prs.slide_masters[0].element
            heading_font, body_font = self._extract_font_scheme(master_element)
        except Exception:
            pass

        # Slide dimensions (EMU)
        margins: dict[str, float] = {}
        try:
            if prs.slide_width is not None and prs.slide_height is not None:
                margins["slide_width_emu"] = float(prs.slide_width)
                margins["slide_height_emu"] = float(prs.slide_height)
        except Exception:
            pass

        # Layout zones — per-layout placeholder positions
        layout_zones = self._extract_layout_zones(prs)

        return DesignProfile(
            primary_color=primary_color,
            secondary_color=secondary_color,
            accent_color=accent_color,
            color_palette=color_palette,
            heading_font=heading_font,
            body_font=body_font,
            available_layouts=available_layouts,
            margins=margins,
            layout_zones=layout_zones,
        )

    # ── DOCX analysis ──────────────────────────────────────────────────────

    async def _analyze_docx(self, path: str) -> DesignProfile:
        """Extract from DOCX: style names, fonts from styles, margins from first section."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx required: pip install firefly-dworkers[document]")
        return await asyncio.to_thread(self._analyze_docx_sync, path)

    def _analyze_docx_sync(self, path: str) -> DesignProfile:
        doc = docx.Document(path)

        # Style names
        styles: list[str] = []
        try:
            styles = [s.name for s in doc.styles if s.type is not None]
        except Exception:
            pass

        # Extract fonts from key styles
        heading_font = ""
        body_font = ""
        font_sizes: dict[str, float] = {}
        try:
            heading_font, body_font, font_sizes = self._extract_docx_fonts(doc)
        except Exception:
            pass

        # Margins from first section
        margins: dict[str, float] = {}
        try:
            if doc.sections:
                section = doc.sections[0]
                for attr in ("left_margin", "right_margin", "top_margin", "bottom_margin"):
                    val = getattr(section, attr, None)
                    if val is not None:
                        # Convert EMU to inches (914400 EMU per inch)
                        margins[attr.replace("_margin", "")] = round(float(val) / 914400.0, 3)
        except Exception:
            pass

        return DesignProfile(
            heading_font=heading_font,
            body_font=body_font,
            font_sizes=font_sizes,
            styles=styles,
            margins=margins,
        )

    # ── XLSX analysis ──────────────────────────────────────────────────────

    async def _analyze_xlsx(self, path: str) -> DesignProfile:
        """Extract from XLSX: default font, column widths, sheet names."""
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl required: pip install firefly-dworkers[data]")
        return await asyncio.to_thread(self._analyze_xlsx_sync, path)

    def _analyze_xlsx_sync(self, path: str) -> DesignProfile:
        wb = openpyxl.load_workbook(path, read_only=True)

        # Default font
        body_font = ""
        font_sizes: dict[str, float] = {}
        try:
            if wb.active is not None:
                cell = wb.active.cell(row=1, column=1)
                if cell.font and cell.font.name:
                    body_font = cell.font.name
                if cell.font and cell.font.sz:
                    font_sizes["default"] = float(cell.font.sz)
        except Exception:
            pass

        # Sheet names as layout info
        available_layouts: list[str] = []
        try:
            available_layouts = wb.sheetnames
        except Exception:
            pass

        wb.close()

        return DesignProfile(
            body_font=body_font,
            font_sizes=font_sizes,
            available_layouts=available_layouts,
        )

    # ── VLM fallback ─────────────────────────────────────────────────────

    @staticmethod
    def _is_profile_empty(profile: DesignProfile) -> bool:
        """Check if the XML-extracted profile has no meaningful design data."""
        return not profile.primary_color and not profile.heading_font and not profile.color_palette

    async def _vlm_analyze(self, source: str, xml_profile: DesignProfile) -> DesignProfile:
        """Render first slide -> VLM -> DesignProfile. Merge with XML results."""
        from firefly_dworkers.design.preview import SlidePreviewRenderer
        from fireflyframework_genai.agents.base import FireflyAgent
        from fireflyframework_genai.types import BinaryContent

        renderer = SlidePreviewRenderer(dpi=150)
        png_list = await renderer.render_presentation(source)
        if not png_list:
            return xml_profile

        agent = FireflyAgent(
            "template-analyzer",
            model=self._vlm_model,
            instructions=(
                "Analyze this slide template. Extract: primary/secondary/accent "
                "colors (hex), heading font, body font, color palette. "
                "Return a DesignProfile."
            ),
            output_type=DesignProfile,
            auto_register=False,
        )
        result = await agent.run([
            "Analyze this template slide and extract the design language.",
            BinaryContent(data=png_list[0], media_type="image/png"),
        ])
        vlm_profile: DesignProfile = result.output

        # Merge: XML non-empty fields take priority
        return DesignProfile(
            primary_color=xml_profile.primary_color or vlm_profile.primary_color,
            secondary_color=xml_profile.secondary_color or vlm_profile.secondary_color,
            accent_color=xml_profile.accent_color or vlm_profile.accent_color,
            background_color=xml_profile.background_color if xml_profile.background_color != "#ffffff" else vlm_profile.background_color,
            text_color=xml_profile.text_color if xml_profile.text_color != "#333333" else vlm_profile.text_color,
            color_palette=xml_profile.color_palette or vlm_profile.color_palette,
            heading_font=xml_profile.heading_font or vlm_profile.heading_font,
            body_font=xml_profile.body_font or vlm_profile.body_font,
            font_sizes=xml_profile.font_sizes or vlm_profile.font_sizes,
            available_layouts=xml_profile.available_layouts or vlm_profile.available_layouts,
            preferred_layouts=xml_profile.preferred_layouts or vlm_profile.preferred_layouts,
            margins=xml_profile.margins or vlm_profile.margins,
            line_spacing=xml_profile.line_spacing if xml_profile.line_spacing != 1.15 else vlm_profile.line_spacing,
            styles=xml_profile.styles or vlm_profile.styles,
            master_slide_names=xml_profile.master_slide_names or vlm_profile.master_slide_names,
            layout_zones=xml_profile.layout_zones or vlm_profile.layout_zones,
        )

    # ── LLM-based placeholder classification ─────────────────────────────

    async def _classify_zones_with_llm(
        self,
        zones: dict[str, LayoutZone],
    ) -> dict[str, LayoutZone]:
        """Use LLM to classify custom placeholders by semantic role.

        Only called when ``self._vlm_model`` is set.  Sends placeholder
        metadata (name, position, size) to the LLM and updates ``ph_type``
        for any placeholders currently marked ``"custom"``.  After
        reclassification, ``title_ph_idx`` and ``body_ph_idx`` are recomputed.
        """
        from pydantic import BaseModel, Field
        from pydantic_ai import Agent

        # ── LLM response schema ──────────────────────────────────────────
        class _PlaceholderRole(BaseModel):
            idx: int
            role: str  # title, subtitle, body, date, client_name, author, logo, decorative, custom

        class _LayoutRoles(BaseModel):
            layout_name: str
            roles: list[_PlaceholderRole] = Field(default_factory=list)

        class _ClassificationResult(BaseModel):
            layouts: list[_LayoutRoles] = Field(default_factory=list)

        # ── Build prompt with all custom placeholders ─────────────────────
        has_custom = False
        descriptions: list[str] = []
        for zone_name, zone in zones.items():
            custom = [ph for ph in zone.placeholders if ph.ph_type == "custom"]
            if not custom:
                continue
            has_custom = True
            desc = f"Layout '{zone_name}':"
            for ph in custom:
                desc += (
                    f"\n  - idx={ph.idx}, name='{ph.name}', "
                    f"position=({ph.left}\", {ph.top}\"), "
                    f"size=({ph.width}\" x {ph.height}\")"
                )
            descriptions.append(desc)

        if not has_custom:
            return zones

        agent: Agent[None, _ClassificationResult] = Agent(
            self._vlm_model,
            output_type=_ClassificationResult,
            system_prompt=(
                "You are a presentation template analyst. Given placeholder "
                "metadata from slide layouts, classify each placeholder's "
                "semantic role. Consider the placeholder name (in any language), "
                "its position on the slide, and its size.\n\n"
                "Valid roles: title, subtitle, body, date, client_name, "
                "author, logo, decorative, custom.\n\n"
                "Rules:\n"
                "- Large text areas near the top are likely 'title' or 'subtitle'\n"
                "- Large text areas in the middle are likely 'body'\n"
                "- Small text areas at the bottom are likely 'date', 'author', "
                "or 'client_name'\n"
                "- Use the placeholder name as a strong hint regardless of language\n"
                "- If unsure, use 'custom'"
            ),
        )

        prompt = (
            "Classify the semantic role of each placeholder below.\n\n"
            + "\n\n".join(descriptions)
        )
        result = await agent.run(prompt)
        classification = result.output

        # ── Apply LLM classifications ─────────────────────────────────────
        role_map: dict[str, dict[int, str]] = {}
        for lr in classification.layouts:
            role_map[lr.layout_name] = {r.idx: r.role for r in lr.roles}

        for zone_name, zone in zones.items():
            layout_roles = role_map.get(zone_name, {})
            changed = False
            for ph in zone.placeholders:
                if ph.ph_type == "custom" and ph.idx in layout_roles:
                    ph.ph_type = layout_roles[ph.idx]
                    changed = True

            if changed:
                # Recompute title_ph_idx and body_ph_idx
                zone.title_ph_idx = self._find_title_ph(zone.placeholders)
                zone.body_ph_idx = self._find_body_ph(
                    zone.placeholders, zone.title_ph_idx
                )

        return zones

    # ── Layout zone extraction ────────────────────────────────────────────

    def _extract_layout_zones(self, prs: Any) -> dict[str, LayoutZone]:
        """Extract placeholder positions for every layout across all masters."""
        zones: dict[str, LayoutZone] = {}
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name in zones:
                    continue
                phs: list[PlaceholderZone] = []
                for ph in layout.placeholders:
                    fmt = ph.placeholder_format
                    phs.append(PlaceholderZone(
                        idx=fmt.idx,
                        name=ph.name,
                        ph_type=self._classify_placeholder(fmt.idx, ph.name),
                        left=round(float(ph.left or 0) / _EMU_PER_INCH, 2),
                        top=round(float(ph.top or 0) / _EMU_PER_INCH, 2),
                        width=round(float(ph.width or 0) / _EMU_PER_INCH, 2),
                        height=round(float(ph.height or 0) / _EMU_PER_INCH, 2),
                    ))
                title_idx = self._find_title_ph(phs)
                body_idx = self._find_body_ph(phs, title_idx)
                cl, ct, cw, ch = self._compute_content_zone(phs, title_idx, prs)
                zones[layout.name] = LayoutZone(
                    layout_name=layout.name,
                    placeholders=phs,
                    content_left=cl,
                    content_top=ct,
                    content_width=cw,
                    content_height=ch,
                    title_ph_idx=title_idx,
                    body_ph_idx=body_idx,
                )
        return zones

    @staticmethod
    def _classify_placeholder(idx: int, name: str) -> str:
        """Classify a placeholder using only OOXML standard indices.

        Non-standard indices are left as ``"custom"`` for LLM classification.
        """
        # OOXML spec-defined placeholder indices — language-independent
        _STANDARD: dict[int, str] = {
            0: "title",
            1: "body",
            10: "date_time",
            11: "slide_number",
            12: "footer",
            15: "title",
        }
        return _STANDARD.get(idx, "custom")

    @staticmethod
    def _find_title_ph(phs: list[PlaceholderZone]) -> int | None:
        """Return idx of the title placeholder (spec-classified only)."""
        for ph in phs:
            if ph.ph_type == "title":
                return ph.idx
        return None

    @staticmethod
    def _find_body_ph(
        phs: list[PlaceholderZone],
        exclude_title_idx: int | None,
    ) -> int | None:
        """Return idx of the largest non-title text placeholder (min 1.5" tall, 3" wide)."""
        best_idx: int | None = None
        best_area = 0.0
        for ph in phs:
            if ph.idx == exclude_title_idx:
                continue
            if ph.ph_type in ("date_time", "slide_number", "footer"):
                continue
            if ph.height < 1.5 or ph.width < 3.0:
                continue
            area = ph.width * ph.height
            if area > best_area:
                best_area = area
                best_idx = ph.idx
        return best_idx

    @staticmethod
    def _compute_content_zone(
        phs: list[PlaceholderZone],
        title_idx: int | None,
        prs: Any,
    ) -> tuple[float, float, float, float]:
        """Compute the safe content area (left, top, width, height) in inches."""
        # If there's a body-sized placeholder, use it as the content zone
        body_ph: PlaceholderZone | None = None
        best_area = 0.0
        for ph in phs:
            if ph.idx == title_idx:
                continue
            if ph.ph_type in ("date_time", "slide_number", "footer"):
                continue
            if ph.height < 1.5 or ph.width < 3.0:
                continue
            area = ph.width * ph.height
            if area > best_area:
                best_area = area
                body_ph = ph
        if body_ph is not None:
            return body_ph.left, body_ph.top, body_ph.width, body_ph.height

        # Fallback: compute from slide dimensions
        slide_w = round(float(prs.slide_width or 0) / _EMU_PER_INCH, 2)
        slide_h = round(float(prs.slide_height or 0) / _EMU_PER_INCH, 2)

        # Find the bottom of the title/subtitle to determine content top
        max_bottom = 0.0
        for ph in phs:
            if ph.ph_type in ("title", "subtitle"):
                bottom = ph.top + ph.height
                if bottom > max_bottom:
                    max_bottom = bottom

        left = 0.5
        top = round(max_bottom + 0.2, 2) if max_bottom > 0 else 1.5
        width = round(slide_w - 1.0, 2) if slide_w > 1.0 else 8.0
        height = round(slide_h - top - 0.5, 2) if slide_h > top + 0.5 else 4.0
        return left, top, width, height

    # ── Private helpers ─────────────────────────────────────────────────────

    @staticmethod
    def _extract_theme_colors(master_element: Any) -> tuple[list[str], str, str, str]:
        """Parse clrScheme from the slide master XML element.

        Returns ``(color_palette, primary_color, secondary_color, accent_color)``.
        """
        color_palette: list[str] = []
        primary_color = ""
        secondary_color = ""
        accent_color = ""

        clr_scheme = master_element.find(f".//{{{_DRAWINGML_NS}}}clrScheme")
        if clr_scheme is None:
            return color_palette, primary_color, secondary_color, accent_color

        # Map of clrScheme child tag local names to semantic meaning
        tag_map: dict[str, str] = {
            "dk1": "dark1",
            "dk2": "dark2",
            "lt1": "light1",
            "lt2": "light2",
            "accent1": "accent1",
            "accent2": "accent2",
            "accent3": "accent3",
            "accent4": "accent4",
            "accent5": "accent5",
            "accent6": "accent6",
            "hlink": "hlink",
            "folHlink": "folHlink",
        }

        for child in clr_scheme:
            local_name = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local_name not in tag_map:
                continue

            color_hex = TemplateAnalyzer._extract_color_value(child)
            if color_hex:
                color_palette.append(color_hex)

                if local_name == "dk1" and not primary_color:
                    primary_color = color_hex
                elif local_name == "dk2" and not secondary_color:
                    secondary_color = color_hex
                elif local_name == "accent1" and not accent_color:
                    accent_color = color_hex

        return color_palette, primary_color, secondary_color, accent_color

    @staticmethod
    def _extract_color_value(element: Element) -> str:
        """Extract a hex color value from a color scheme child element.

        Looks for ``srgbClr`` (val attr) and ``sysClr`` (lastClr attr) sub-elements.
        """
        for sub in element:
            local = sub.tag.split("}")[-1] if "}" in sub.tag else sub.tag
            if local == "srgbClr":
                val = sub.get("val", "")
                if val:
                    return f"#{val}"
            elif local == "sysClr":
                last_clr = sub.get("lastClr", "")
                if last_clr:
                    return f"#{last_clr}"
        return ""

    @staticmethod
    def _extract_font_scheme(master_element: Any) -> tuple[str, str]:
        """Parse fontScheme from the slide master XML element.

        Returns ``(heading_font, body_font)``.
        """
        heading_font = ""
        body_font = ""

        font_scheme = master_element.find(f".//{{{_DRAWINGML_NS}}}fontScheme")
        if font_scheme is None:
            return heading_font, body_font

        # Major font (headings)
        major = font_scheme.find(f"{{{_DRAWINGML_NS}}}majorFont")
        if major is not None:
            latin = major.find(f"{{{_DRAWINGML_NS}}}latin")
            if latin is not None:
                heading_font = latin.get("typeface", "")

        # Minor font (body)
        minor = font_scheme.find(f"{{{_DRAWINGML_NS}}}minorFont")
        if minor is not None:
            latin = minor.find(f"{{{_DRAWINGML_NS}}}latin")
            if latin is not None:
                body_font = latin.get("typeface", "")

        return heading_font, body_font

    @staticmethod
    def _extract_docx_fonts(doc: Any) -> tuple[str, str, dict[str, float]]:
        """Extract font info from DOCX styles.

        Returns ``(heading_font, body_font, font_sizes)``.
        """
        heading_font = ""
        body_font = ""
        font_sizes: dict[str, float] = {}

        for style in doc.styles:
            if style.type is None:
                continue
            try:
                name = style.name
                font = style.font
                if font is None:
                    continue

                if name == "Normal" and font.name:
                    body_font = font.name
                    if font.size:
                        font_sizes["body"] = round(float(font.size) / 12700.0, 1)
                elif name == "Heading 1" and font.name:
                    heading_font = font.name
                    if font.size:
                        font_sizes["h1"] = round(float(font.size) / 12700.0, 1)
                elif name == "Heading 2" and font.name and font.size:
                    font_sizes["h2"] = round(float(font.size) / 12700.0, 1)
            except Exception:
                continue

        return heading_font, body_font, font_sizes
