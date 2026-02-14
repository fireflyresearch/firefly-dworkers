"""VLM-based visual refinement loop for presentation positioning.

Renders PPTX slides to PNG, sends to a VLM to identify positioning issues,
applies rule-based fixes to shape positions, and iterates until acceptable.
"""

from __future__ import annotations

import io
import logging
from typing import Any

from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

_EMU_PER_INCH = 914400


# ── Data models ────────────────────────────────────────────────────────────


class PositionFix(BaseModel):
    """A single positioning fix to apply to a shape."""

    slide_index: int
    shape_name: str = ""
    issue: str = ""
    fix_type: str = ""  # move, resize, font_reduce
    new_left: float | None = None  # EMU
    new_top: float | None = None
    new_width: float | None = None
    new_height: float | None = None


class RefinementFeedback(BaseModel):
    """VLM feedback for a single slide."""

    slide_index: int = 0
    score: float = Field(ge=0, le=10, default=5.0)
    position_fixes: list[PositionFix] = Field(default_factory=list)
    is_acceptable: bool = True


# ── System prompt ─────────────────────────────────────────────────────────

_REFINEMENT_SYSTEM_PROMPT = """\
You are a presentation layout quality inspector. Analyze this slide image \
for POSITIONING issues only. Focus on:

1. Content overlapping decorative/header elements
2. Text overflow outside placeholder bounds
3. Charts or tables placed outside the content area
4. Elements too close to slide edges (< 0.3 inches)
5. Overlapping elements

For each issue found, provide:
- slide_index: the slide number (0-based)
- shape_name: name of the shape to fix (if identifiable)
- issue: brief description
- fix_type: one of "move", "resize", "font_reduce"
- new_left, new_top, new_width, new_height: corrected position in EMU \
  (1 inch = 914400 EMU). Only provide fields that need changing.

Score the slide 0-10 for positioning quality. Set is_acceptable=True if \
score >= 7.
"""


# ── VisualRefiner ─────────────────────────────────────────────────────────


class VisualRefiner:
    """Iteratively refines PPTX positioning using VLM feedback."""

    def __init__(
        self,
        *,
        vlm_model: str,
        max_iterations: int = 2,
        score_threshold: float = 7.0,
    ) -> None:
        self._vlm_model = vlm_model
        self._max_iterations = max_iterations
        self._score_threshold = score_threshold

    async def refine(
        self,
        pptx_bytes: bytes,
        layout_zones: dict[str, Any] | None = None,
    ) -> bytes:
        """Iteratively refine PPTX positioning using VLM feedback.

        Returns the (possibly improved) PPTX bytes.
        """
        current = pptx_bytes

        for iteration in range(self._max_iterations):
            # 1. Render all slides to PNG
            try:
                png_list = await self._render_slides(current)
            except Exception:
                logger.warning("Refinement: render failed, returning current PPTX")
                return current

            if not png_list:
                return current

            # 2. Send each slide to VLM for feedback
            all_fixes: list[PositionFix] = []
            all_acceptable = True

            for slide_idx, png_bytes in enumerate(png_list):
                try:
                    feedback = await self._evaluate_slide(slide_idx, png_bytes)
                except Exception:
                    logger.warning(
                        "Refinement: VLM evaluation failed for slide %d",
                        slide_idx,
                    )
                    continue

                if not feedback.is_acceptable:
                    all_acceptable = False
                if feedback.position_fixes:
                    all_fixes.extend(feedback.position_fixes)

            # 3. If all slides acceptable, we're done
            if all_acceptable or not all_fixes:
                logger.info(
                    "Refinement: all slides acceptable after %d iteration(s)",
                    iteration + 1,
                )
                return current

            # 4. Apply fixes
            try:
                current = self._apply_fixes(current, all_fixes)
            except Exception:
                logger.warning("Refinement: fix application failed, returning current")
                return current

        return current

    async def _render_slides(self, pptx_bytes: bytes) -> list[bytes]:
        """Render PPTX bytes to list of PNG bytes."""
        import tempfile

        from firefly_dworkers.design.preview import SlidePreviewRenderer

        renderer = SlidePreviewRenderer(dpi=150)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(pptx_bytes)
            tmp_path = tmp.name

        import os

        try:
            return await renderer.render_presentation(tmp_path)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    async def _evaluate_slide(
        self,
        slide_index: int,
        png_bytes: bytes,
    ) -> RefinementFeedback:
        """Send a slide PNG to the VLM and get positioning feedback."""
        from fireflyframework_genai.agents.base import FireflyAgent
        from fireflyframework_genai.types import BinaryContent

        agent = FireflyAgent(
            "refinement-inspector",
            model=self._vlm_model,
            instructions=_REFINEMENT_SYSTEM_PROMPT,
            output_type=RefinementFeedback,
            auto_register=False,
        )

        result = await agent.run([
            f"Evaluate slide {slide_index} for positioning quality.",
            BinaryContent(data=png_bytes, media_type="image/png"),
        ])
        feedback: RefinementFeedback = result.output
        feedback.slide_index = slide_index
        return feedback

    @staticmethod
    def _apply_fixes(pptx_bytes: bytes, fixes: list[PositionFix]) -> bytes:
        """Apply position fixes to PPTX shapes by name or slide index."""
        import pptx as pptx_mod

        prs = pptx_mod.Presentation(io.BytesIO(pptx_bytes))

        for fix in fixes:
            if fix.slide_index >= len(prs.slides):
                continue
            slide = prs.slides[fix.slide_index]

            # Find shape by name
            target = None
            if fix.shape_name:
                for shape in slide.shapes:
                    if shape.name == fix.shape_name:
                        target = shape
                        break

            if target is None:
                continue

            if fix.new_left is not None:
                target.left = int(fix.new_left)
            if fix.new_top is not None:
                target.top = int(fix.new_top)
            if fix.new_width is not None:
                target.width = int(fix.new_width)
            if fix.new_height is not None:
                target.height = int(fix.new_height)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
