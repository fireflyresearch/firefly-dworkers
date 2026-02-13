"""OS-agnostic slide preview renderer using matplotlib.

Renders visual approximations of PPTX slides as PNG bytes — no OS-specific
tools required (no qlmanage, Keynote, or LibreOffice). Designed to produce
~80% visual fidelity previews suitable for VLM-based quality evaluation.
"""

from __future__ import annotations

import io
import logging
from typing import Any

logger = logging.getLogger(__name__)

# EMU (English Metric Units) conversion: 914400 EMU = 1 inch
_EMU_PER_INCH = 914400


def _emu_to_inches(emu: int | float) -> float:
    return float(emu) / _EMU_PER_INCH


class SlidePreviewRenderer:
    """Renders PPTX slides to PNG bytes using matplotlib."""

    def __init__(self, dpi: int = 150) -> None:
        self._dpi = dpi

    async def render_slide(
        self,
        slide: Any,
        slide_width: int | float,
        slide_height: int | float,
    ) -> bytes:
        """Render a single python-pptx slide to PNG bytes."""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as patches

        w_in = _emu_to_inches(slide_width)
        h_in = _emu_to_inches(slide_height)

        fig, ax = plt.subplots(figsize=(w_in, h_in), dpi=self._dpi)
        try:
            # White background
            ax.set_xlim(0, w_in)
            ax.set_ylim(0, h_in)
            ax.invert_yaxis()
            ax.set_facecolor("white")
            ax.set_aspect("equal")
            ax.axis("off")
            fig.patch.set_facecolor("white")

            for shape in slide.shapes:
                self._draw_shape(ax, shape, w_in, h_in)

            buf = io.BytesIO()
            fig.savefig(buf, format="png", bbox_inches="tight", dpi=self._dpi,
                        facecolor="white", edgecolor="none")
            return buf.getvalue()
        finally:
            plt.close(fig)

    async def render_presentation(self, pptx_path: str) -> list[bytes]:
        """Render all slides in a PPTX file to a list of PNG bytes."""
        import pptx as pptx_mod

        prs = pptx_mod.Presentation(pptx_path)
        results: list[bytes] = []
        for slide in prs.slides:
            png = await self.render_slide(slide, prs.slide_width, prs.slide_height)
            results.append(png)
        return results

    def _draw_shape(self, ax: Any, shape: Any, slide_w: float, slide_h: float) -> None:
        """Dispatch shape drawing based on type."""
        import matplotlib.patches as patches

        try:
            x = _emu_to_inches(shape.left) if shape.left else 0
            y = _emu_to_inches(shape.top) if shape.top else 0
            w = _emu_to_inches(shape.width) if shape.width else 0
            h = _emu_to_inches(shape.height) if shape.height else 0
        except Exception:
            return

        if w <= 0 or h <= 0:
            return

        if shape.has_chart:
            self._draw_chart_placeholder(ax, x, y, w, h, shape)
        elif shape.has_table:
            self._draw_table(ax, x, y, w, h, shape)
        elif shape.has_text_frame:
            self._draw_text_frame(ax, x, y, w, h, shape)
        elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            self._draw_image_placeholder(ax, x, y, w, h)

    @staticmethod
    def _draw_text_frame(
        ax: Any, x: float, y: float, w: float, h: float, shape: Any
    ) -> None:
        """Draw text at the shape's position."""
        text = shape.text_frame.text.strip()
        if not text:
            return
        # Truncate long text
        display = text[:120] + "..." if len(text) > 120 else text
        # Estimate font size from shape height
        font_size = min(max(h * 6, 6), 18)
        ax.text(
            x + 0.1, y + h / 2, display,
            fontsize=font_size, va="center", ha="left",
            color="#333333", clip_on=True,
            fontfamily="sans-serif",
        )

    @staticmethod
    def _draw_chart_placeholder(
        ax: Any, x: float, y: float, w: float, h: float, shape: Any
    ) -> None:
        """Draw a colored rectangle representing a chart."""
        import matplotlib.patches as patches

        rect = patches.FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.05",
            facecolor="#E8F0FE", edgecolor="#4285F4",
            linewidth=1.5,
        )
        ax.add_patch(rect)
        title = ""
        try:
            if shape.chart.has_title:
                title = shape.chart.chart_title.text_frame.text
        except Exception:
            pass
        label = title or "[Chart]"
        ax.text(
            x + w / 2, y + h / 2, label,
            fontsize=9, va="center", ha="center",
            color="#4285F4", fontweight="bold",
        )

    @staticmethod
    def _draw_table(
        ax: Any, x: float, y: float, w: float, h: float, shape: Any
    ) -> None:
        """Draw a grid of rectangles approximating a table."""
        import matplotlib.patches as patches

        table = shape.table
        n_rows = len(table.rows)
        n_cols = len(table.columns)
        if n_rows == 0 or n_cols == 0:
            return

        cell_w = w / n_cols
        cell_h = h / n_rows

        for r_idx, row in enumerate(table.rows):
            for c_idx in range(n_cols):
                cx = x + c_idx * cell_w
                cy = y + r_idx * cell_h
                # Header row gets darker background
                fc = "#D6E4F0" if r_idx == 0 else ("#F5F5F5" if r_idx % 2 == 0 else "white")
                rect = patches.Rectangle(
                    (cx, cy), cell_w, cell_h,
                    facecolor=fc, edgecolor="#CCCCCC", linewidth=0.5,
                )
                ax.add_patch(rect)
                try:
                    text = table.cell(r_idx, c_idx).text.strip()
                    if text:
                        display = text[:15] + ".." if len(text) > 15 else text
                        fs = 6 if r_idx > 0 else 7
                        ax.text(
                            cx + cell_w / 2, cy + cell_h / 2, display,
                            fontsize=fs, va="center", ha="center",
                            color="#333333", clip_on=True,
                            fontweight="bold" if r_idx == 0 else "normal",
                        )
                except Exception:
                    pass

    @staticmethod
    def _draw_image_placeholder(
        ax: Any, x: float, y: float, w: float, h: float
    ) -> None:
        """Draw a dashed rectangle placeholder for images."""
        import matplotlib.patches as patches

        rect = patches.Rectangle(
            (x, y), w, h,
            facecolor="#F9F9F9", edgecolor="#AAAAAA",
            linewidth=1, linestyle="--",
        )
        ax.add_patch(rect)
        ax.text(
            x + w / 2, y + h / 2, "[Image]",
            fontsize=8, va="center", ha="center", color="#AAAAAA",
        )
