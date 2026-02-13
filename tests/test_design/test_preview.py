"""Tests for SlidePreviewRenderer -- OS-agnostic slide-to-PNG rendering."""

from __future__ import annotations

import tempfile

import pytest


class TestSlidePreviewRenderer:
    """Tests for the matplotlib-based slide preview renderer."""

    async def test_render_empty_slide(self) -> None:
        """Blank slide produces valid PNG bytes."""
        pptx_mod = pytest.importorskip("pptx")
        pytest.importorskip("matplotlib")
        from firefly_dworkers.design.preview import SlidePreviewRenderer

        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        renderer = SlidePreviewRenderer(dpi=72)
        png = await renderer.render_slide(slide, prs.slide_width, prs.slide_height)
        assert len(png) > 0
        assert png[:4] == b"\x89PNG"

    async def test_render_slide_with_text(self) -> None:
        """Slide with text shapes produces valid PNG."""
        pptx_mod = pytest.importorskip("pptx")
        pytest.importorskip("matplotlib")
        from firefly_dworkers.design.preview import SlidePreviewRenderer

        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide.shapes.title.text = "Hello World"
        renderer = SlidePreviewRenderer(dpi=72)
        png = await renderer.render_slide(slide, prs.slide_width, prs.slide_height)
        assert png[:4] == b"\x89PNG"

    async def test_render_slide_with_table(self) -> None:
        """Slide with a table produces valid PNG."""
        pptx_mod = pytest.importorskip("pptx")
        pytest.importorskip("matplotlib")
        from pptx.util import Inches

        from firefly_dworkers.design.preview import SlidePreviewRenderer

        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1))
        table_shape.table.cell(0, 0).text = "Header"
        table_shape.table.cell(1, 0).text = "Value"
        renderer = SlidePreviewRenderer(dpi=72)
        png = await renderer.render_slide(slide, prs.slide_width, prs.slide_height)
        assert png[:4] == b"\x89PNG"

    async def test_render_presentation_multiple(self) -> None:
        """List length matches slide count."""
        pptx_mod = pytest.importorskip("pptx")
        pytest.importorskip("matplotlib")
        from firefly_dworkers.design.preview import SlidePreviewRenderer

        prs = pptx_mod.Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            tmp_path = f.name

        renderer = SlidePreviewRenderer(dpi=72)
        result = await renderer.render_presentation(tmp_path)
        assert len(result) == 3
        for png in result:
            assert png[:4] == b"\x89PNG"

    async def test_render_slide_with_chart(self) -> None:
        """Slide with a chart shape produces valid PNG with chart placeholder."""
        pptx_mod = pytest.importorskip("pptx")
        pytest.importorskip("matplotlib")
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Emu

        from firefly_dworkers.design.preview import SlidePreviewRenderer

        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("S1", [10, 20])
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(914400), Emu(914400), Emu(5000000), Emu(3000000),
            chart_data,
        )
        renderer = SlidePreviewRenderer(dpi=72)
        png = await renderer.render_slide(slide, prs.slide_width, prs.slide_height)
        assert png[:4] == b"\x89PNG"
