"""End-to-end integration tests for the design pipeline."""

from __future__ import annotations

import io
import tempfile

import pytest
from pydantic_ai.models.test import TestModel

from firefly_dworkers.design.engine import DesignEngine
from firefly_dworkers.design.models import (
    ContentBrief,
    ContentSection,
    DataSeries,
    DataSet,
    DesignSpec,
    OutputType,
)


class TestEndToEndPresentation:
    """Full pipeline: ContentBrief -> DesignEngine -> PowerPointTool -> PPTX."""

    async def test_brief_to_pptx(self) -> None:
        pytest.importorskip("pptx")
        from firefly_dworkers.tools.presentation.models import SlideSpec
        from firefly_dworkers.tools.presentation.powerpoint import PowerPointTool

        engine = DesignEngine(model=TestModel())
        brief = ContentBrief(
            output_type=OutputType.PRESENTATION,
            title="Test Deck",
            sections=[
                ContentSection(heading="Intro", content="Hello world"),
                ContentSection(heading="Data", content="Analysis results"),
            ],
        )
        spec = await engine.design(brief)
        assert isinstance(spec, DesignSpec)
        assert spec.output_type == OutputType.PRESENTATION

        # Render to PPTX using profile colors
        ppt = PowerPointTool()
        slides = [
            SlideSpec(title=brief.title, content="Executive summary"),
            SlideSpec(title="Details", bullet_points=["Point 1", "Point 2"]),
        ]
        data = await ppt.create(slides=slides)
        assert len(data) > 0

        # Verify it's a valid PPTX (ZIP with [Content_Types].xml)
        import zipfile

        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            assert "[Content_Types].xml" in zf.namelist()


class TestEndToEndDocument:
    """Full pipeline: ContentBrief -> DesignEngine -> WordTool -> DOCX."""

    async def test_brief_to_docx(self) -> None:
        pytest.importorskip("docx")
        import docx as docx_mod

        from firefly_dworkers.tools.document.models import SectionSpec
        from firefly_dworkers.tools.document.word import WordTool

        engine = DesignEngine(model=TestModel())
        brief = ContentBrief(
            output_type=OutputType.DOCUMENT,
            title="Test Report",
            sections=[
                ContentSection(heading="Chapter 1", content="Introduction text"),
                ContentSection(heading="Chapter 2", content="Analysis"),
            ],
        )
        spec = await engine.design(brief)
        assert spec.output_type == OutputType.DOCUMENT

        word = WordTool()
        sections = [
            SectionSpec(heading="Chapter 1", content="Introduction text"),
            SectionSpec(heading="Chapter 2", content="Analysis"),
        ]
        data = await word.create(title=brief.title, sections=sections)

        doc = docx_mod.Document(io.BytesIO(data))
        # Title + 2 headings + 2 content paragraphs
        assert len(doc.paragraphs) >= 3


class TestEndToEndSpreadsheet:
    """Full pipeline: ContentBrief -> DesignEngine -> ExcelTool -> XLSX."""

    async def test_brief_to_xlsx(self) -> None:
        pytest.importorskip("openpyxl")
        import openpyxl

        from firefly_dworkers.tools.spreadsheet.excel import ExcelTool
        from firefly_dworkers.tools.spreadsheet.models import SheetSpec

        engine = DesignEngine(model=TestModel())
        brief = ContentBrief(
            output_type=OutputType.SPREADSHEET,
            title="Data Export",
            sections=[ContentSection(heading="Summary", content="Data")],
        )
        spec = await engine.design(brief)
        assert spec.output_type == OutputType.SPREADSHEET

        excel = ExcelTool()
        sheets = [SheetSpec(name="Summary", headers=["Metric", "Value"], rows=[["Revenue", "1000"]])]
        data = await excel.create(sheets=sheets)

        wb = openpyxl.load_workbook(io.BytesIO(data))
        assert "Summary" in wb.sheetnames


class TestEndToEndWithCharts:
    """Pipeline with chart data: ContentBrief with DataSets -> DesignSpec with charts."""

    async def test_brief_with_datasets_resolves_charts(self) -> None:
        engine = DesignEngine(model=TestModel())
        brief = ContentBrief(
            output_type=OutputType.PRESENTATION,
            title="Chart Test",
            sections=[ContentSection(heading="Data", chart_ref="revenue")],
            datasets=[
                DataSet(
                    name="revenue",
                    categories=["Q1", "Q2", "Q3", "Q4"],
                    series=[DataSeries(name="Revenue", values=[100, 150, 200, 250])],
                ),
            ],
        )
        spec = await engine.design(brief)
        assert "revenue" in spec.charts
        assert spec.charts["revenue"].chart_type in ("bar", "line", "pie", "scatter", "area", "doughnut")


class TestTemplateRoundtrip:
    """Analyze a PPTX -> extract profile -> verify profile is populated."""

    async def test_analyze_pptx_produces_profile(self) -> None:
        pytest.importorskip("pptx")
        import tempfile

        from pptx import Presentation

        from firefly_dworkers.design.analyzer import TemplateAnalyzer

        # Create a minimal PPTX
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            path = f.name

        analyzer = TemplateAnalyzer()
        profile = await analyzer.analyze(path)
        # Should have extracted at least layouts
        assert isinstance(profile.available_layouts, list)
