"""Tests for PowerPointTool adapter."""

from __future__ import annotations

import io
import os
import struct
import tempfile
import zlib

import pytest
from fireflyframework_genai.tools.base import BaseTool

from firefly_dworkers.design.models import TextStyle
from firefly_dworkers.tools.presentation.base import PresentationTool
from firefly_dworkers.tools.presentation.models import ChartSpec, SlideSpec, TableSpec
from firefly_dworkers.tools.presentation.powerpoint import PowerPointTool
from firefly_dworkers.tools.registry import tool_registry


def _make_minimal_png(width: int = 1, height: int = 1) -> bytes:
    """Create a minimal valid 1x1 white PNG image in memory."""

    def _chunk(chunk_type: bytes, data: bytes) -> bytes:
        c = chunk_type + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    header = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    ihdr = _chunk(b"IHDR", ihdr_data)
    # Single white pixel row: filter byte (0) + RGB (255,255,255)
    raw_row = b"\x00" + b"\xff\xff\xff" * width
    raw_data = raw_row * height
    idat = _chunk(b"IDAT", zlib.compress(raw_data))
    iend = _chunk(b"IEND", b"")
    return header + ihdr + idat + iend


class TestPowerPointToolRegistration:
    def test_is_presentation_tool(self) -> None:
        assert issubclass(PowerPointTool, PresentationTool)

    def test_is_base_tool(self) -> None:
        assert issubclass(PowerPointTool, BaseTool)

    def test_registry_entry(self) -> None:
        assert tool_registry.has("powerpoint")
        assert tool_registry.get_class("powerpoint") is PowerPointTool

    def test_category(self) -> None:
        assert tool_registry.get_category("powerpoint") == "presentation"

    def test_name(self) -> None:
        assert PowerPointTool().name == "powerpoint"


class TestPowerPointToolRead:
    async def test_read_presentation(self) -> None:
        pptx = pytest.importorskip("pptx")

        # Create a minimal .pptx in memory
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Title"
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(buf.read())
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            result = await tool.execute(action="read", source=tmp_path)
            assert "slides" in result
            assert len(result["slides"]) == 1
            assert result["slides"][0]["title"] == "Test Title"
        finally:
            os.unlink(tmp_path)


class TestPowerPointToolCreate:
    async def test_create_presentation_basic(self) -> None:
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(title="Slide 1", content="Hello world").model_dump(),
            SlideSpec(title="Slide 2", bullet_points=["A", "B", "C"]).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        assert result["bytes_length"] > 0

    async def test_create_with_table(self) -> None:
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Data Table",
                table={"headers": ["Name", "Value"], "rows": [["A", "1"]]},
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True


class TestPowerPointToolPublicAPI:
    async def test_create_returns_bytes(self) -> None:
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        result = await tool.create(slides=[SlideSpec(title="Slide 1", content="Hello")])
        assert isinstance(result, bytes)
        assert len(result) > 0

    async def test_create_and_save(self, tmp_path) -> None:
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        out = str(tmp_path / "test.pptx")
        path = await tool.create_and_save(out, slides=[SlideSpec(title="Test")])
        assert os.path.exists(path)
        assert os.path.getsize(path) > 0

    async def test_artifact_bytes_after_execute(self) -> None:
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        await tool.execute(action="create", slides=[SlideSpec(title="T").model_dump()])
        assert tool.artifact_bytes is not None
        assert len(tool.artifact_bytes) > 0


class TestPowerPointToolCharts:
    async def test_create_with_chart(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Revenue Chart",
                chart=ChartSpec(
                    chart_type="bar",
                    title="Revenue by Quarter",
                    categories=["Q1", "Q2", "Q3"],
                    series=[{"name": "2025", "values": [100, 200, 300]}],
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        # Verify chart shape exists
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
        assert len(chart_shapes) >= 1

    async def test_create_with_chart_title(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Chart Slide",
                chart=ChartSpec(
                    chart_type="line",
                    title="Growth Trend",
                    categories=["Jan", "Feb", "Mar"],
                    series=[{"name": "Revenue", "values": [10, 20, 30]}],
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
        assert chart_shapes[0].chart.chart_title.text_frame.text == "Growth Trend"

    async def test_create_with_pie_chart(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Market Share",
                chart=ChartSpec(
                    chart_type="pie",
                    title="Market Share",
                    categories=["Product A", "Product B", "Product C"],
                    series=[{"name": "Share", "values": [45, 30, 25]}],
                    show_legend=True,
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
        assert len(chart_shapes) == 1

    async def test_create_with_multi_series_chart(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Comparison",
                chart=ChartSpec(
                    chart_type="bar",
                    title="Year over Year",
                    categories=["Q1", "Q2", "Q3", "Q4"],
                    series=[
                        {"name": "2024", "values": [100, 150, 200, 250]},
                        {"name": "2025", "values": [120, 180, 220, 300]},
                    ],
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
        assert len(chart_shapes) == 1

    async def test_no_chart_when_spec_is_none(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [SlideSpec(title="No Chart", content="Just text").model_dump()]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
        assert len(chart_shapes) == 0


class TestPowerPointToolImages:
    async def test_create_with_image(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        png_data = _make_minimal_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png_data)
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            slides = [
                SlideSpec(title="Image Slide", image_path=tmp_path).model_dump(),
            ]
            result = await tool.execute(action="create", slides=slides)
            assert result["success"] is True
            data = tool.artifact_bytes
            assert data is not None
            prs = pptx_mod.Presentation(io.BytesIO(data))
            # Find picture shapes (shape_type 13 = MSO_SHAPE_TYPE.PICTURE)
            pic_shapes = [
                s for s in prs.slides[0].shapes if s.shape_type == 13  # noqa: PLR2004
            ]
            assert len(pic_shapes) >= 1
        finally:
            os.unlink(tmp_path)

    async def test_create_with_image_placement(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        from firefly_dworkers.design.models import ImagePlacement

        png_data = _make_minimal_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png_data)
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            slides = [
                SlideSpec(
                    title="Placed Images",
                    images=[
                        ImagePlacement(
                            file_path=tmp_path,
                            left=914400,
                            top=1524000,
                            width=3048000,
                            height=2286000,
                        ),
                    ],
                ).model_dump(),
            ]
            result = await tool.execute(action="create", slides=slides)
            assert result["success"] is True
            data = tool.artifact_bytes
            assert data is not None
            prs = pptx_mod.Presentation(io.BytesIO(data))
            pic_shapes = [
                s for s in prs.slides[0].shapes if s.shape_type == 13  # noqa: PLR2004
            ]
            assert len(pic_shapes) >= 1
        finally:
            os.unlink(tmp_path)

    async def test_create_with_both_image_path_and_images(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        from firefly_dworkers.design.models import ImagePlacement

        png_data = _make_minimal_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png_data)
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            slides = [
                SlideSpec(
                    title="Multiple Images",
                    image_path=tmp_path,
                    images=[
                        ImagePlacement(file_path=tmp_path, left=500000, top=500000),
                    ],
                ).model_dump(),
            ]
            result = await tool.execute(action="create", slides=slides)
            assert result["success"] is True
            data = tool.artifact_bytes
            assert data is not None
            prs = pptx_mod.Presentation(io.BytesIO(data))
            pic_shapes = [
                s for s in prs.slides[0].shapes if s.shape_type == 13  # noqa: PLR2004
            ]
            # Should have at least 2: one from image_path, one from images list
            assert len(pic_shapes) >= 2  # noqa: PLR2004
        finally:
            os.unlink(tmp_path)

    async def test_empty_images_list_is_noop(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(title="No Images", content="text", images=[]).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        pic_shapes = [
            s for s in prs.slides[0].shapes if s.shape_type == 13  # noqa: PLR2004
        ]
        assert len(pic_shapes) == 0


class TestPowerPointToolStyling:
    async def test_create_with_title_style(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Styled Title",
                content="Some content",
                title_style=TextStyle(
                    font_name="Arial",
                    font_size=28,
                    bold=True,
                    color="#1a73e8",
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        title_shape = prs.slides[0].shapes.title
        run = title_shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Arial"
        assert run.font.bold is True

    async def test_create_with_body_style(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Styled Body",
                content="Body text here",
                body_style=TextStyle(
                    font_name="Calibri",
                    font_size=16,
                    italic=True,
                    color="#333333",
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        # Find body placeholder
        body_ph = None
        for shape in prs.slides[0].placeholders:
            if shape.placeholder_format.idx == 1:
                body_ph = shape
                break
        assert body_ph is not None
        run = body_ph.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Calibri"
        assert run.font.italic is True

    async def test_create_with_bullet_points_and_body_style(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Styled Bullets",
                bullet_points=["Point A", "Point B"],
                body_style=TextStyle(
                    font_name="Helvetica",
                    font_size=14,
                    bold=True,
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        body_ph = None
        for shape in prs.slides[0].placeholders:
            if shape.placeholder_format.idx == 1:
                body_ph = shape
                break
        assert body_ph is not None
        run = body_ph.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Helvetica"
        assert run.font.bold is True

    async def test_title_style_color_applied(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        from pptx.dml.color import RGBColor

        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Color Title",
                content="text",
                title_style=TextStyle(color="#ff5500"),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        title_shape = prs.slides[0].shapes.title
        run = title_shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(0xFF, 0x55, 0x00)

    async def test_no_style_is_noop(self) -> None:
        """When no style is provided, the tool should still succeed."""
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(title="Plain Title", content="Plain content").model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True

    async def test_apply_text_style_with_none(self) -> None:
        """_apply_text_style should be a no-op when style is None."""
        pytest.importorskip("pptx")
        # Just verify it does not raise
        PowerPointTool._apply_text_style(None, None)


class TestPowerPointToolCombined:
    async def test_chart_and_styling_together(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Full Featured Slide",
                content="Some body text",
                title_style=TextStyle(font_name="Georgia", font_size=32, bold=True),
                body_style=TextStyle(font_name="Verdana", font_size=14),
                chart=ChartSpec(
                    chart_type="bar",
                    title="Sales",
                    categories=["A", "B"],
                    series=[{"name": "S1", "values": [10, 20]}],
                ),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        slide = prs.slides[0]
        # Title styled
        title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert title_run.font.name == "Georgia"
        # Chart present
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) >= 1

    async def test_image_and_chart_together(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        png_data = _make_minimal_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png_data)
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            slides = [
                SlideSpec(
                    title="Chart + Image",
                    image_path=tmp_path,
                    chart=ChartSpec(
                        chart_type="line",
                        title="Trend",
                        categories=["X", "Y"],
                        series=[{"name": "S1", "values": [5, 10]}],
                    ),
                ).model_dump(),
            ]
            result = await tool.execute(action="create", slides=slides)
            assert result["success"] is True
            data = tool.artifact_bytes
            assert data is not None
            prs = pptx_mod.Presentation(io.BytesIO(data))
            slide = prs.slides[0]
            chart_shapes = [s for s in slide.shapes if s.has_chart]
            pic_shapes = [s for s in slide.shapes if s.shape_type == 13]  # noqa: PLR2004
            assert len(chart_shapes) >= 1
            assert len(pic_shapes) >= 1
        finally:
            os.unlink(tmp_path)

    async def test_backwards_compatible_basic_create(self) -> None:
        """Existing basic create functionality still works unchanged."""
        pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [
            SlideSpec(title="Slide 1", content="Hello world").model_dump(),
            SlideSpec(title="Slide 2", bullet_points=["A", "B", "C"]).model_dump(),
            SlideSpec(
                title="Data",
                table={"headers": ["Name", "Value"], "rows": [["X", "1"]]},
                speaker_notes="Remember this.",
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        assert result["bytes_length"] > 0


class TestPowerPointTextAlignment:
    """Tests for text alignment support."""

    async def test_text_alignment_center(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        from pptx.enum.text import PP_ALIGN

        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Centered",
                content="Center me",
                body_style=TextStyle(alignment="center"),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        body_ph = None
        for shape in prs.slides[0].placeholders:
            if shape.placeholder_format.idx == 1:
                body_ph = shape
                break
        assert body_ph is not None
        assert body_ph.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER

    async def test_text_alignment_default_left(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        from pptx.enum.text import PP_ALIGN

        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Left Aligned",
                content="Default left",
                body_style=TextStyle(alignment="left"),
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        body_ph = None
        for shape in prs.slides[0].placeholders:
            if shape.placeholder_format.idx == 1:
                body_ph = shape
                break
        assert body_ph is not None
        assert body_ph.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT

    async def test_bullet_spacing(self) -> None:
        pptx_mod = pytest.importorskip("pptx")
        from pptx.util import Pt

        tool = PowerPointTool()
        slides = [
            SlideSpec(
                title="Bullets",
                bullet_points=["Point A", "Point B"],
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        data = tool.artifact_bytes
        assert data is not None
        prs = pptx_mod.Presentation(io.BytesIO(data))
        body_ph = None
        for shape in prs.slides[0].placeholders:
            if shape.placeholder_format.idx == 1:
                body_ph = shape
                break
        assert body_ph is not None
        p = body_ph.text_frame.paragraphs[0]
        assert p.space_before == Pt(4)
        assert p.space_after == Pt(2)


class TestPowerPointTableStyling:
    """Tests for consulting-quality table styling."""

    def _make_table_slide(self, table_spec):
        """Helper: create a single-slide presentation with a table and return the table."""
        pptx_mod = pytest.importorskip("pptx")
        prs = pptx_mod.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        PowerPointTool._add_table(slide, table_spec)
        table_shapes = [s for s in slide.shapes if s.has_table]
        assert len(table_shapes) == 1
        return table_shapes[0].table

    def test_table_header_background(self) -> None:
        """Verify header cells have solid fill XML when header_bg_color set."""
        table = self._make_table_slide(
            TableSpec(
                headers=["A", "B"],
                rows=[["1", "2"]],
                header_bg_color="#1a3c6d",
            )
        )
        header_cell = table.cell(0, 0)
        xml = header_cell._tc.xml
        assert "solidFill" in xml
        assert "1A3C6D" in xml

    def test_table_alternating_rows(self) -> None:
        """Verify even/odd data rows differ when alternating_rows=True."""
        table = self._make_table_slide(
            TableSpec(
                headers=["Col"],
                rows=[["r0"], ["r1"], ["r2"], ["r3"]],
                alternating_rows=True,
                alt_row_color="#F5F5F5",
            )
        )
        # Row 1 (index 0, even) should have fill, row 2 (index 1, odd) should not
        even_xml = table.cell(1, 0)._tc.xml  # data row 0
        odd_xml = table.cell(2, 0)._tc.xml  # data row 1
        assert "F5F5F5" in even_xml
        assert "F5F5F5" not in odd_xml

    def test_table_borders_present(self) -> None:
        """Verify cell XML contains border elements."""
        table = self._make_table_slide(
            TableSpec(
                headers=["X"],
                rows=[["y"]],
                border_color="#CCCCCC",
            )
        )
        xml = table.cell(0, 0)._tc.xml
        assert "lnT" in xml
        assert "lnB" in xml
        assert "lnL" in xml
        assert "lnR" in xml

    def test_table_default_styling(self) -> None:
        """Minimal TableSpec with no styling overrides still works."""
        table = self._make_table_slide(
            TableSpec(headers=["Name", "Value"], rows=[["A", "1"]])
        )
        assert table.cell(0, 0).text == "Name"
        assert table.cell(1, 0).text == "A"

    def test_table_dict_input_backward_compat(self) -> None:
        """Raw dict input (legacy API) still works."""
        table = self._make_table_slide(
            {"headers": ["H1", "H2"], "rows": [["v1", "v2"]]}
        )
        assert table.cell(0, 0).text == "H1"
        assert table.cell(1, 1).text == "v2"

    def test_table_cell_margins_set(self) -> None:
        """Verify cell padding margins are set in XML."""
        table = self._make_table_slide(
            TableSpec(headers=["A"], rows=[["b"]])
        )
        xml = table.cell(0, 0)._tc.xml
        assert "marL" in xml
        assert "marT" in xml


# ── Zone-aware placement tests ──────────────────────────────────────────

from firefly_dworkers.tools.presentation.models import ContentZone


def _make_content_zone(
    left_inches: float = 0.5,
    top_inches: float = 1.5,
    width_inches: float = 9.0,
    height_inches: float = 5.0,
    title_ph_idx: int | None = 0,
    body_ph_idx: int | None = 1,
    **extra_map: int,
) -> ContentZone:
    """Helper to create a ContentZone in EMU from inch values."""
    EMU = 914400
    ph_map: dict[str, int] = {}
    if title_ph_idx is not None:
        ph_map["title"] = title_ph_idx
    if body_ph_idx is not None:
        ph_map["body"] = body_ph_idx
    ph_map.update(extra_map)
    return ContentZone(
        left=int(left_inches * EMU),
        top=int(top_inches * EMU),
        width=int(width_inches * EMU),
        height=int(height_inches * EMU),
        title_ph_idx=title_ph_idx,
        body_ph_idx=body_ph_idx,
        placeholder_map=ph_map,
    )


class TestZoneAwareTitlePlacement:
    async def test_title_uses_zone_title_ph_idx(self) -> None:
        """When content_zone has title_ph_idx, title is written to that placeholder."""
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        # Default pptx has ph[0] as title — zone also says ph[0]
        cz = _make_content_zone(title_ph_idx=0)
        slides = [SlideSpec(title="Zone Title", content="body", content_zone=cz).model_dump()]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
        # Title should be filled
        assert prs.slides[0].shapes.title.text == "Zone Title"

    async def test_title_fallback_without_zone(self) -> None:
        """Without content_zone, falls back to slide.shapes.title."""
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [SlideSpec(title="Fallback Title", content="body").model_dump()]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
        assert prs.slides[0].shapes.title.text == "Fallback Title"


class TestZoneAwareBodyPlacement:
    async def test_body_uses_zone_body_ph_idx(self) -> None:
        """When content_zone has body_ph_idx, content goes to that placeholder."""
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        cz = _make_content_zone(body_ph_idx=1)
        slides = [SlideSpec(title="T", content="Zone body", content_zone=cz).model_dump()]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
        body_ph = None
        for shape in prs.slides[0].placeholders:
            if shape.placeholder_format.idx == 1:
                body_ph = shape
                break
        assert body_ph is not None
        assert body_ph.text_frame.text == "Zone body"

    async def test_body_fallback_largest_area(self) -> None:
        """Without zone, body picks largest non-title text placeholder."""
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        slides = [SlideSpec(title="T", content="Fallback body").model_dump()]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
        # Should still find a body placeholder via fallback
        body_found = False
        for shape in prs.slides[0].placeholders:
            if shape.has_text_frame and shape.text_frame.text == "Fallback body":
                body_found = True
                break
        assert body_found


class TestZoneAwareChartPlacement:
    async def test_chart_positioned_in_content_zone(self) -> None:
        """When content_zone is set, chart should use zone coordinates."""
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        cz = _make_content_zone(left_inches=1.5, top_inches=2.0, width_inches=7.0, height_inches=4.0)
        slides = [
            SlideSpec(
                title="Chart",
                chart=ChartSpec(
                    chart_type="bar",
                    title="Sales",
                    categories=["A", "B"],
                    series=[{"name": "S1", "values": [10, 20]}],
                ),
                content_zone=cz,
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
        chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
        assert len(chart_shapes) >= 1
        chart_shape = chart_shapes[0]
        # Chart should be positioned at the zone coordinates
        EMU = 914400
        assert chart_shape.left == int(1.5 * EMU)
        assert chart_shape.top == int(2.0 * EMU)


class TestZoneAwareTablePlacement:
    async def test_table_positioned_in_content_zone(self) -> None:
        """When content_zone is set, table should use zone left/top/width."""
        pptx_mod = pytest.importorskip("pptx")
        tool = PowerPointTool()
        cz = _make_content_zone(left_inches=1.5, top_inches=2.0, width_inches=7.0, height_inches=4.0)
        slides = [
            SlideSpec(
                title="Table",
                table=TableSpec(headers=["A", "B"], rows=[["1", "2"]]),
                content_zone=cz,
            ).model_dump(),
        ]
        result = await tool.execute(action="create", slides=slides)
        assert result["success"] is True
        prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
        table_shapes = [s for s in prs.slides[0].shapes if s.has_table]
        assert len(table_shapes) >= 1
        tbl_shape = table_shapes[0]
        EMU = 914400
        assert tbl_shape.left == int(1.5 * EMU)
        assert tbl_shape.top == int(2.0 * EMU)


class TestZoneAwareImagePlacement:
    async def test_image_positioned_in_content_zone(self) -> None:
        """When content_zone is set, image should be within zone bounds."""
        pptx_mod = pytest.importorskip("pptx")
        png_data = _make_minimal_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png_data)
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            cz = _make_content_zone(left_inches=1.5, top_inches=2.0, width_inches=7.0, height_inches=4.0)
            slides = [
                SlideSpec(title="Image", image_path=tmp_path, content_zone=cz).model_dump(),
            ]
            result = await tool.execute(action="create", slides=slides)
            assert result["success"] is True
            prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
            pic_shapes = [s for s in prs.slides[0].shapes if s.shape_type == 13]  # noqa: PLR2004
            assert len(pic_shapes) >= 1
            EMU = 914400
            pic = pic_shapes[0]
            assert pic.left == int(1.5 * EMU)
            assert pic.top == int(2.0 * EMU)
        finally:
            os.unlink(tmp_path)


class TestZoneFallbackPreservesDefaults:
    async def test_all_positions_fallback_without_zone(self) -> None:
        """Without content_zone, all elements use existing hardcoded defaults."""
        pptx_mod = pytest.importorskip("pptx")
        from pptx.util import Inches

        png_data = _make_minimal_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(png_data)
            tmp_path = f.name

        try:
            tool = PowerPointTool()
            slides = [
                SlideSpec(
                    title="Defaults",
                    content="text",
                    table=TableSpec(headers=["H"], rows=[["r"]]),
                    image_path=tmp_path,
                ).model_dump(),
            ]
            result = await tool.execute(action="create", slides=slides)
            assert result["success"] is True
            prs = pptx_mod.Presentation(io.BytesIO(tool.artifact_bytes))
            slide = prs.slides[0]

            # Table at default position: Inches(1), Inches(2)
            table_shapes = [s for s in slide.shapes if s.has_table]
            assert len(table_shapes) >= 1
            assert table_shapes[0].left == Inches(1)
            assert table_shapes[0].top == Inches(2)

            # Image at default position: Inches(1), Inches(2)
            pic_shapes = [s for s in slide.shapes if s.shape_type == 13]  # noqa: PLR2004
            assert len(pic_shapes) >= 1
            assert pic_shapes[0].left == Inches(1)
            assert pic_shapes[0].top == Inches(2)
        finally:
            os.unlink(tmp_path)
